#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
第01课《数字的奥秘》教学课件生成器
卡通风格、活泼设计、炫酷动画
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_UNDERLINE

# 卡通配色方案 - 紫蓝科技风
COLORS = {
    'primary': RGBColor(102, 0, 153),      # 清华紫
    'secondary': RGBColor(148, 7, 10),     # 北大红
    'accent': RGBColor(255, 107, 107),     # 珊瑚红
    'bg_light': RGBColor(240, 248, 255),   # 浅蓝背景
    'success': RGBColor(76, 175, 80),      # 成功绿
    'warning': RGBColor(255, 193, 7),      # 警告黄
    'info': RGBColor(33, 150, 243),        # 信息蓝
    'white': RGBColor(255, 255, 255),
    'black': RGBColor(33, 33, 33),
    'gray': RGBColor(158, 158, 158),
}

def create_presentation():
    """创建演示文稿"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs

def add_background(slide, color):
    """添加纯色背景"""
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = color
    background.line.fill.background()
    # 将背景移到最底层
    spTree = slide.shapes._spTree
    sp = background._element
    spTree.remove(sp)
    spTree.insert(2, sp)

def add_title_box(slide, title, subtitle=None, top=Inches(0.5)):
    """添加标题框"""
    # 标题背景
    title_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), top,
        Inches(12.333), Inches(1.2)
    )
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = COLORS['primary']
    title_box.line.fill.background()
    
    # 标题文字
    title_tf = title_box.text_frame
    title_tf.text = title
    title_p = title_tf.paragraphs[0]
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = COLORS['white']
    title_p.font.name = 'Microsoft YaHei'
    title_p.alignment = PP_ALIGN.CENTER
    title_tf.word_wrap = True
    
    if subtitle:
        # 副标题
        sub_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), top + Inches(1.3),
            Inches(12.333), Inches(0.6)
        )
        sub_box.fill.solid()
        sub_box.fill.fore_color.rgb = COLORS['secondary']
        sub_box.line.fill.background()
        
        sub_tf = sub_box.text_frame
        sub_tf.text = subtitle
        sub_p = sub_tf.paragraphs[0]
        sub_p.font.size = Pt(24)
        sub_p.font.color.rgb = COLORS['white']
        sub_p.font.name = 'Microsoft YaHei'
        sub_p.alignment = PP_ALIGN.CENTER

def add_content_box(slide, left, top, width, height, title, content, title_color):
    """添加内容框"""
    # 外框
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    box.fill.solid()
    box.fill.fore_color.rgb = COLORS['white']
    box.line.color.rgb = title_color
    box.line.width = Pt(3)
    
    # 标题条
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, Inches(0.6)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = title_color
    title_bar.line.fill.background()
    
    # 标题文字
    title_tf = title_bar.text_frame
    title_tf.text = title
    title_p = title_tf.paragraphs[0]
    title_p.font.size = Pt(20)
    title_p.font.bold = True
    title_p.font.color.rgb = COLORS['white']
    title_p.font.name = 'Microsoft YaHei'
    title_p.alignment = PP_ALIGN.CENTER
    
    # 内容文字
    content_box = slide.shapes.add_textbox(
        left + Inches(0.2), top + Inches(0.8),
        width - Inches(0.4), height - Inches(1)
    )
    content_tf = content_box.text_frame
    content_tf.text = content
    content_p = content_tf.paragraphs[0]
    content_p.font.size = Pt(18)
    content_p.font.color.rgb = COLORS['black']
    content_p.font.name = 'Microsoft YaHei'
    content_p.line_spacing = 1.5
    content_tf.word_wrap = True

def add_emoji_decorations(slide):
    """添加emoji装饰"""
    emojis = ['🔢', '✨', '🎯', '🎮']
    positions = [
        (Inches(0.3), Inches(0.3)),
        (Inches(12.5), Inches(0.3)),
        (Inches(0.3), Inches(6.8)),
        (Inches(12.5), Inches(6.8)),
    ]
    
    for emoji, (left, top) in zip(emojis, positions):
        emoji_box = slide.shapes.add_textbox(left, top, Inches(0.5), Inches(0.5))
        emoji_tf = emoji_box.text_frame
        emoji_tf.text = emoji
        emoji_p = emoji_tf.paragraphs[0]
        emoji_p.font.size = Pt(36)

def create_slide_1_cover(prs):
    """第1页：封面页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    add_background(slide, COLORS['bg_light'])
    
    # 主标题
    title_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(2),
        Inches(11.333), Inches(1.5)
    )
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = COLORS['primary']
    title_box.line.fill.background()
    
    title_tf = title_box.text_frame
    title_tf.text = '🔢 第01课《数字的奥秘》'
    title_p = title_tf.paragraphs[0]
    title_p.font.size = Pt(54)
    title_p.font.bold = True
    title_p.font.color.rgb = COLORS['white']
    title_p.font.name = 'Microsoft YaHei'
    title_p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    sub_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2), Inches(4),
        Inches(9.333), Inches(0.8)
    )
    sub_box.fill.solid()
    sub_box.fill.fore_color.rgb = COLORS['secondary']
    sub_box.line.fill.background()
    
    sub_tf = sub_box.text_frame
    sub_tf.text = '🎯 发现规律 · 培养数感 · 快乐学习'
    sub_p = sub_tf.paragraphs[0]
    sub_p.font.size = Pt(28)
    sub_p.font.color.rgb = COLORS['white']
    sub_p.font.name = 'Microsoft YaHei'
    sub_p.alignment = PP_ALIGN.CENTER
    
    # 底部装饰
    deco_box = slide.shapes.add_textbox(
        Inches(4), Inches(5.5),
        Inches(5.333), Inches(1)
    )
    deco_tf = deco_box.text_frame
    deco_tf.text = '🐸 🔢 ✨ 🎯 🎮 📚 💡'
    deco_p = deco_tf.paragraphs[0]
    deco_p.font.size = Pt(48)
    deco_p.alignment = PP_ALIGN.CENTER
    
    add_emoji_decorations(slide)
    return slide

def create_slide_2_objectives(prs):
    """第2页：课程目标"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS['bg_light'])
    add_title_box(slide, '🎯 本课学习目标', 'Learning Objectives')
    
    objectives = [
        ('📊 认识数字规律', '学习等差数列、递推数列、斐波那契数列'),
        ('🧩 掌握数独入门', '理解3×3数独规则，运用排除法解题'),
        ('🎮 培养数学兴趣', '通过游戏化学习，爱上数学思考'),
        ('💡 提升思维能力', '观察力、逻辑力、计算力、创造力'),
    ]
    
    for i, (title, desc) in enumerate(objectives):
        row = i // 2
        col = i % 2
        left = Inches(0.8 + col * 6.2)
        top = Inches(2.2 + row * 2.5)
        
        add_content_box(slide, left, top, Inches(5.8), Inches(2.2),
                       title, desc, COLORS['info'])
    
    add_emoji_decorations(slide)
    return slide

def create_slide_3_warmup(prs):
    """第3页：热身游戏"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS['bg_light'])
    add_title_box(slide, '🎮 热身游戏：数字接龙', 'Warm-up Game')
    
    # 游戏规则
    rule_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(2),
        Inches(6), Inches(4.5)
    )
    rule_box.fill.solid()
    rule_box.fill.fore_color.rgb = COLORS['white']
    rule_box.line.color.rgb = COLORS['success']
    rule_box.line.width = Pt(3)
    
    rule_title = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(2),
        Inches(6), Inches(0.7)
    )
    rule_title.fill.solid()
    rule_title.fill.fore_color.rgb = COLORS['success']
    rule_title.line.fill.background()
    
    rule_tf = rule_title.text_frame
    rule_tf.text = '📋 游戏规则'
    rule_p = rule_tf.paragraphs[0]
    rule_p.font.size = Pt(24)
    rule_p.font.bold = True
    rule_p.font.color.rgb = COLORS['white']
    rule_p.alignment = PP_ALIGN.CENTER
    
    # 规则内容
    content_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(2.9),
        Inches(5.6), Inches(3.5)
    )
    content_tf = content_box.text_frame
    content_tf.text = '''1️⃣ 按规律接龙，接不上表演节目

2️⃣ 从1数到30，遇3的倍数拍手
   1, 2, 👏, 4, 5, 👏...

3️⃣ 发现规律：3, 6, 9, 12, 15...'''
    content_p = content_tf.paragraphs[0]
    content_p.font.size = Pt(20)
    content_p.font.name = 'Microsoft YaHei'
    content_p.line_spacing = 1.6
    
    # 示例数列
    example_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.8), Inches(2),
        Inches(6), Inches(4.5)
    )
    example_box.fill.solid()
    example_box.fill.fore_color.rgb = COLORS['white']
    example_box.line.color.rgb = COLORS['warning']
    example_box.line.width = Pt(3)
    
    example_title = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.8), Inches(2),
        Inches(6), Inches(0.7)
    )
    example_title.fill.solid()
    example_title.fill.fore_color.rgb = COLORS['warning']
    example_title.line.fill.background()
    
    example_tf = example_title.text_frame
    example_tf.text = '🔢 数字规律示例'
    example_p = example_tf.paragraphs[0]
    example_p.font.size = Pt(24)
    example_p.font.bold = True
    example_p.font.color.rgb = COLORS['white']
    example_p.alignment = PP_ALIGN.CENTER
    
    # 示例内容
    example_content = slide.shapes.add_textbox(
        Inches(7), Inches(2.9),
        Inches(5.6), Inches(3.5)
    )
    example_tf = example_content.text_frame
    example_tf.text = '''① 2, 4, 6, 8, __, __
   规律：每次+2

② 1, 3, 5, 7, __, __
   规律：奇数序列

③ 1, 2, 4, 7, 11, __
   规律：差递增'''
    example_p = example_tf.paragraphs[0]
    example_p.font.size = Pt(20)
    example_p.font.name = 'Microsoft YaHei'
    example_p.line_spacing = 1.6
    
    add_emoji_decorations(slide)
    return slide

def create_slide_4_concept1(prs):
    """第4页：等差数列概念"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS['bg_light'])
    add_title_box(slide, '📚 概念一：等差数列', 'Arithmetic Sequence')
    
    # 定义框
    def_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(2),
        Inches(12.333), Inches(1.5)
    )
    def_box.fill.solid()
    def_box.fill.fore_color.rgb = COLORS['info']
    def_box.line.fill.background()
    
    def_tf = def_box.text_frame
    def_tf.text = '📖 定义：相邻两个数的差都相等的数列叫做等差数列'
    def_p = def_tf.paragraphs[0]
    def_p.font.size = Pt(26)
    def_p.font.bold = True
    def_p.font.color.rgb = COLORS['white']
    def_p.alignment = PP_ALIGN.CENTER
    
    # 示例
    examples = [
        ('🐸 小青蛙跳荷叶', '1, 3, 5, 7, 9, 11...', '每次+2'),
        ('🚂 小火车车厢号', '1, 4, 7, 10, 13...', '每次+3'),
        ('🐿️ 小松鼠藏松果', '2, 4, 6, 8, 10...', '每次+2'),
    ]
    
    for i, (title, sequence, rule) in enumerate(examples):
        left = Inches(0.5 + i * 4.2)
        top = Inches(4)
        
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, Inches(4), Inches(2.8)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = COLORS['white']
        box.line.color.rgb = COLORS['primary']
        box.line.width = Pt(2)
        
        # 标题
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, Inches(4), Inches(0.6)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = COLORS['primary']
        title_bar.line.fill.background()
        
        title_tf = title_bar.text_frame
        title_tf.text = title
        title_p = title_tf.paragraphs[0]
        title_p.font.size = Pt(18)
        title_p.font.bold = True
        title_p.font.color.rgb = COLORS['white']
        title_p.alignment = PP_ALIGN.CENTER
        
        # 内容
        content_box = slide.shapes.add_textbox(
            left + Inches(0.2), top + Inches(0.8),
            Inches(3.6), Inches(1.8)
        )
        content_tf = content_box.text_frame
        content_tf.text = f'{sequence}\n\n规律：{rule}'
        content_p = content_tf.paragraphs[0]
        content_p.font.size = Pt(18)
        content_p.font.name = 'Microsoft YaHei'
        content_p.line_spacing = 1.5
    
    add_emoji_decorations(slide)
    return slide

def create_slide_5_concept2(prs):
    """第5页：斐波那契数列"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS['bg_light'])
    add_title_box(slide, '📚 概念二：斐波那契数列', 'Fibonacci Sequence')
    
    # 定义框
    def_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(2),
        Inches(12.333), Inches(1.3)
    )
    def_box.fill.solid()
    def_box.fill.fore_color.rgb = COLORS['accent']
    def_box.line.fill.background()
    
    def_tf = def_box.text_frame
    def_tf.text = '🐰 兔子数列：从第3个数开始，每个数 = 前两个数之和'
    def_p = def_tf.paragraphs[0]
    def_p.font.size = Pt(24)
    def_p.font.bold = True
    def_p.font.color.rgb = COLORS['white']
    def_p.alignment = PP_ALIGN.CENTER
    
    # 数列展示
    seq_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(3.5),
        Inches(11.333), Inches(1.2)
    )
    seq_box.fill.solid()
    seq_box.fill.fore_color.rgb = COLORS['white']
    seq_box.line.color.rgb = COLORS['warning']
    seq_box.line.width = Pt(3)
    
    seq_tf = seq_box.text_frame
    seq_tf.text = '1 , 1 , 2 , 3 , 5 , 8 , 13 , 21 , 34 , 55 ...'
    seq_p = seq_tf.paragraphs[0]
    seq_p.font.size = Pt(32)
    seq_p.font.name = 'Courier New'
    seq_p.font.bold = True
    seq_p.font.color.rgb = COLORS['primary']
    seq_p.alignment = PP_ALIGN.CENTER
    
    # 计算过程
    calc_box = slide.shapes.add_textbox(
        Inches(1), Inches(5),
        Inches(11.333), Inches(2)
    )
    calc_tf = calc_box.text_frame
    calc_tf.text = '''计算过程：
1 + 1 = 2    1 + 2 = 3    2 + 3 = 5    3 + 5 = 8    5 + 8 = 13 ...

✨ 自然界中的斐波那契：花瓣数、松果螺旋、向日葵种子排列'''
    calc_p = calc_tf.paragraphs[0]
    calc_p.font.size = Pt(20)
    calc_p.font.name = 'Microsoft YaHei'
    calc_p.line_spacing = 1.5
    
    add_emoji_decorations(slide)
    return slide

def create_slide_6_example1(prs):
    """第6页：例题1"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS['bg_light'])
    add_title_box(slide, '📝 例题1：找规律填数', 'Example 1')
    
    # 题目框
    q_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(2),
        Inches(12.333), Inches(2)
    )
    q_box.fill.solid()
    q_box.fill.fore_color.rgb = COLORS['white']
    q_box.line.color.rgb = COLORS['info']
    q_box.line.width = Pt(3)
    
    q_tf = q_box.text_frame
    q_tf.text = '🐸 小青蛙跳荷叶\n它踩到的荷叶编号是：1, 3, 5, 7, ___, ___\n请问：接下来会踩到哪两片荷叶？'
    q_p = q_tf.paragraphs[0]
    q_p.font.size = Pt(24)
    q_p.font.name = 'Microsoft YaHei'
    q_p.line_spacing = 1.5
    
    # 解答框
    a_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(4.3),
        Inches(12.333), Inches(2.5)
    )
    a_box.fill.solid()
    a_box.fill.fore_color.rgb = COLORS['success']
    a_box.line.fill.background()
    
    a_tf = a_box.text_frame
    a_tf.text = '''✅ 解答：

观察：3-1=2, 5-3=2, 7-5=2  →  这是一个等差数列，每次+2

答案：7+2=9, 9+2=11

所以接下来会踩到 9 和 11 号荷叶'''
    a_p = a_tf.paragraphs[0]
    a_p.font.size = Pt(22)
    a_p.font.name = 'Microsoft YaHei'
    a_p.font.color.rgb = COLORS['white']
    a_p.line_spacing = 1.4
    
    add_emoji_decorations(slide)
    return slide

def create_slide_7_example2(prs):
    """第7页：例题2 - 数独"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS['bg_light'])
    add_title_box(slide, '📝 例题2：图形数独', 'Example 2 - Sudoku')
    
    # 题目
    q_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(2),
        Inches(6), Inches(4.5)
    )
    q_box.fill.solid()
    q_box.fill.fore_color.rgb = COLORS['white']
    q_box.line.color.rgb = COLORS['warning']
    q_box.line.width = Pt(3)
    
    q_title = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(2),
        Inches(6), Inches(0.6)
    )
    q_title.fill.solid()
    q_title.fill.fore_color.rgb = COLORS['warning']
    q_title.line.fill.background()
    
    q_tf = q_title.text_frame
    q_tf.text = '🧩 题目'
    q_p = q_tf.paragraphs[0]
    q_p.font.size = Pt(22)
    q_p.font.bold = True
    q_p.font.color.rgb = COLORS['white']
    q_p.alignment = PP_ALIGN.CENTER
    
    # 数独题目内容
    q_content = slide.shapes.add_textbox(
        Inches(0.7), Inches(2.8),
        Inches(5.6), Inches(3.5)
    )
    q_tf = q_content.text_frame
    q_tf.text = '''将 🌸、🍀、⭐ 填入空格
使每行每列都有这三种图案

🌸 🍀 ？
⭐ 🌸 ？
？  ？  ？'''
    q_p = q_tf.paragraphs[0]
    q_p.font.size = Pt(28)
    q_p.font.name = 'Microsoft YaHei'
    q_p.line_spacing = 1.6
    
    # 解答
    a_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7), Inches(2),
        Inches(5.833), Inches(4.5)
    )
    a_box.fill.solid()
    a_box.fill.fore_color.rgb = COLORS['success']
    a_box.line.fill.background()
    
    a_title = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7), Inches(2),
        Inches(5.833), Inches(0.6)
    )
    a_title.fill.solid()
    a_title.fill.fore_color.rgb = COLORS['white']
    a_title.line.fill.background()
    
    a_tf = a_title.text_frame
    a_tf.text = '✅ 答案'
    a_p = a_tf.paragraphs[0]
    a_p.font.size = Pt(22)
    a_p.font.bold = True
    a_p.font.color.rgb = COLORS['success']
    a_p.alignment = PP_ALIGN.CENTER
    
    # 答案内容
    a_content = slide.shapes.add_textbox(
        Inches(7.2), Inches(2.8),
        Inches(5.4), Inches(3.5)
    )
    a_tf = a_content.text_frame
    a_tf.text = '''🌸 🍀 ⭐
⭐ 🌸 🍀
🍀 ⭐ 🌸

方法：排除法
第一行缺⭐ → 填⭐
第二行缺🍀 → 填🍀
以此类推...'''
    a_p = a_tf.paragraphs[0]
    a_p.font.size = Pt(24)
    a_p.font.name = 'Microsoft YaHei'
    a_p.font.color.rgb = COLORS['white']
    a_p.line_spacing = 1.6
    
    add_emoji_decorations(slide)
    return slide

def create_slide_8_practice(prs):
    """第8页：课堂练习"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS['bg_light'])
    add_title_box(slide, '🎮 课堂练习', 'Practice Time')
    
    practices = [
        ('练习1 ⭐', '2, 4, 6, 8, __, __\n规律：每次+2\n答案：10, 12'),
        ('练习2 ⭐⭐', '1, 1, 2, 3, 5, 8, __\n规律：斐波那契\n答案：13'),
        ('练习3 ⭐⭐⭐', '🔴+🔴=8, 🔴+🔵=10\n求：🔵+🔵=?\n答案：12'),
    ]
    
    for i, (title, content) in enumerate(practices):
        left = Inches(0.5 + i * 4.2)
        top = Inches(2.2)
        
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, Inches(4), Inches(4.5)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = COLORS['white']
        box.line.color.rgb = COLORS['primary']
        box.line.width = Pt(2)
        
        # 标题条
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, Inches(4), Inches(0.6)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = COLORS['primary']
        title_bar.line.fill.background()
        
        title_tf = title_bar.text_frame
        title_tf.text = title
        title_p = title_tf.paragraphs[0]
        title_p.font.size = Pt(20)
        title_p.font.bold = True
        title_p.font.color.rgb = COLORS['white']
        title_p.alignment = PP_ALIGN.CENTER
        
        # 内容
        content_box = slide.shapes.add_textbox(
            left + Inches(0.2), top + Inches(0.8),
            Inches(3.6), Inches(3.5)
        )
        content_tf = content_box.text_frame
        content_tf.text = content
        content_p = content_tf.paragraphs[0]
        content_p.font.size = Pt(18)
        content_p.font.name = 'Microsoft YaHei'
        content_p.line_spacing = 1.5
    
    add_emoji_decorations(slide)
    return slide

def create_slide_9_summary(prs):
    """第9页：知识总结"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS['bg_light'])
    add_title_box(slide, '🌟 知识总结', 'Summary')
    
    summaries = [
        ('📊 等差数列', '相邻数差相等\n如：2, 4, 6, 8...'),
        ('🐰 斐波那契', '前两项之和\n如：1, 1, 2, 3, 5...'),
        ('🧩 数独规则', '每行每列不重复\n用排除法解题'),
        ('💡 解题方法', '先观察找规律\n再计算验证'),
    ]
    
    for i, (title, content) in enumerate(summaries):
        row = i // 2
        col = i % 2
        left = Inches(0.8 + col * 6.2)
        top = Inches(2.2 + row * 2.5)
        
        colors = [COLORS['info'], COLORS['accent'], COLORS['success'], COLORS['warning']]
        
        add_content_box(slide, left, top, Inches(5.8), Inches(2.2),
                       title, content, colors[i])
    
    add_emoji_decorations(slide)
    return slide

def create_slide_10_homework(prs):
    """第10页：作业布置"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS['bg_light'])
    add_title_box(slide, '📝 作业布置', 'Homework')
    
    homeworks = [
        ('基础作业 ⭐', '完成练习册第1-3题\n（约10分钟）', COLORS['success']),
        ('提高作业 ⭐⭐', '找一找生活中的数字规律\n和家长分享（约15分钟）', COLORS['warning']),
        ('挑战作业 ⭐⭐⭐', '尝试完成一道4×4数独\n（选做）', COLORS['accent']),
    ]
    
    for i, (title, content, color) in enumerate(homeworks):
        left = Inches(0.5 + i * 4.2)
        top = Inches(2.5)
        
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, Inches(4), Inches(3.5)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = COLORS['white']
        box.line.color.rgb = color
        box.line.width = Pt(3)
        
        # 标题条
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, Inches(4), Inches(0.7)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = color
        title_bar.line.fill.background()
        
        title_tf = title_bar.text_frame
        title_tf.text = title
        title_p = title_tf.paragraphs[0]
        title_p.font.size = Pt(22)
        title_p.font.bold = True
        title_p.font.color.rgb = COLORS['white']
        title_p.alignment = PP_ALIGN.CENTER
        
        # 内容
        content_box = slide.shapes.add_textbox(
            left + Inches(0.2), top + Inches(0.9),
            Inches(3.6), Inches(2.4)
        )
        content_tf = content_box.text_frame
        content_tf.text = content
        content_p = content_tf.paragraphs[0]
        content_p.font.size = Pt(18)
        content_p.font.name = 'Microsoft YaHei'
        content_p.line_spacing = 1.5
    
    # 提示
    tip_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(6.2),
        Inches(11.333), Inches(0.8)
    )
    tip_box.fill.solid()
    tip_box.fill.fore_color.rgb = COLORS['info']
    tip_box.line.fill.background()
    
    tip_tf = tip_box.text_frame
    tip_tf.text = '💡 温馨提示：遇到困难可以请教家长，或扫码咨询夏老师'
    tip_p = tip_tf.paragraphs[0]
    tip_p.font.size = Pt(20)
    tip_p.font.color.rgb = COLORS['white']
    tip_p.alignment = PP_ALIGN.CENTER
    
    add_emoji_decorations(slide)
    return slide

def create_slide_11_ending(prs):
    """第11页：结束页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS['primary'])
    
    # 主标题
    title_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(2.5),
        Inches(11.333), Inches(1.5)
    )
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = COLORS['white']
    title_box.line.fill.background()
    
    title_tf = title_box.text_frame
    title_tf.text = '🎉 今天的课程结束啦！'
    title_p = title_tf.paragraphs[0]
    title_p.font.size = Pt(48)
    title_p.font.bold = True
    title_p.font.color.rgb = COLORS['primary']
    title_p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    sub_box = slide.shapes.add_textbox(
        Inches(1), Inches(4.5),
        Inches(11.333), Inches(1)
    )
    sub_tf = sub_box.text_frame
    sub_tf.text = '下节课预告：《图形变变变》🎨'
    sub_p = sub_tf.paragraphs[0]
    sub_p.font.size = Pt(32)
    sub_p.font.color.rgb = COLORS['white']
    sub_p.alignment = PP_ALIGN.CENTER
    
    # 装饰
    deco_box = slide.shapes.add_textbox(
        Inches(4), Inches(5.8),
        Inches(5.333), Inches(1)
    )
    deco_tf = deco_box.text_frame
    deco_tf.text = '✨ 🔢 🎯 🎮 📚 💡 ✨'
    deco_p = deco_tf.paragraphs[0]
    deco_p.font.size = Pt(40)
    deco_p.alignment = PP_ALIGN.CENTER
    
    # 底部信息
    footer = slide.shapes.add_textbox(
        Inches(1), Inches(6.8),
        Inches(11.333), Inches(0.5)
    )
    footer_tf = footer.text_frame
    footer_tf.text = '二年级思维训练课程 · 夏老师原创'
    footer_p = footer_tf.paragraphs[0]
    footer_p.font.size = Pt(16)
    footer_p.font.color.rgb = COLORS['white']
    footer_p.alignment = PP_ALIGN.CENTER
    
    return slide

# 主程序
if __name__ == '__main__':
    print('🎨 正在生成第01课《数字的奥秘》教学课件...')
    
    prs = create_presentation()
    
    # 创建11页幻灯片
    create_slide_1_cover(prs)
    print('✅ 第1页：封面页')
    
    create_slide_2_objectives(prs)
    print('✅ 第2页：课程目标')
    
    create_slide_3_warmup(prs)
    print('✅ 第3页：热身游戏')
    
    create_slide_4_concept1(prs)
    print('✅ 第4页：等差数列概念')
    
    create_slide_5_concept2(prs)
    print('✅ 第5页：斐波那契数列')
    
    create_slide_6_example1(prs)
    print('✅ 第6页：例题1')
    
    create_slide_7_example2(prs)
    print('✅ 第7页：例题2')
    
    create_slide_8_practice(prs)
    print('✅ 第8页：课堂练习')
    
    create_slide_9_summary(prs)
    print('✅ 第9页：知识总结')
    
    create_slide_10_homework(prs)
    print('✅ 第10页：作业布置')
    
    create_slide_11_ending(prs)
    print('✅ 第11页：结束页')
    
    # 保存文件
    output_path = '第01课_数字的奥秘_教学课件.pptx'
    prs.save(output_path)
    print(f'\n🎉 课件生成完成！')
    print(f'📁 文件保存位置：{output_path}')
    print('\n💡 使用说明：')
    print('   1. 在PowerPoint中打开课件')
    print('   2. 按Alt+F11打开VBA编辑器')
    print('   3. 导入"第01课_PPT动画宏代码.bas"')
    print('   4. 运行宏添加动画效果')
